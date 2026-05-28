#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Convertir ESTACION-5-SLIDES.md a PowerPoint PPTX

Requiere: pip install python-pptx
Uso: python convert_slides_to_pptx.py
"""

import re
import sys
import os
from pathlib import Path

# Configurar UTF-8 en Windows
if sys.platform == 'win32':
    os.environ['PYTHONIOENCODING'] = 'utf-8'

# Detectar si python-pptx está disponible
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def parse_markdown_slides(markdown_content):
    """Parsea el markdown y extrae slides separadas por ---"""
    slides = []

    # Dividir por separador de slides (----)
    parts = re.split(r'\n---\n', markdown_content)

    for part in parts:
        if part.strip():
            slides.append(part.strip())

    return slides


def extract_slide_title_and_content(slide_text):
    """Extrae título y contenido de un slide"""
    lines = slide_text.split('\n')
    title = ""
    content = []
    in_code = False

    for i, line in enumerate(lines):
        # Primera línea con # es el título
        if not title and line.startswith('#'):
            title = line.lstrip('#').strip()
        elif line.startswith('```'):
            in_code = not in_code
            if in_code:
                content.append(line)
        else:
            content.append(line)

    return title, '\n'.join(content)


def create_pptx_presentation(slides_data):
    """Crea una presentación PowerPoint desde los datos de slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Colores personalizados
    COLOR_TITLE = RGBColor(0, 51, 102)      # Azul oscuro
    COLOR_SUBTITLE = RGBColor(102, 102, 102) # Gris
    COLOR_TEXT = RGBColor(51, 51, 51)       # Gris oscuro

    for i, slide_text in enumerate(slides_data, 1):
        # Crear slide con layout en blanco
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Extraer título y contenido
        title, content = extract_slide_title_and_content(slide_text)

        # Agregar fondo blanco
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Agregar título
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.word_wrap = True

            # Formatear título
            for paragraph in title_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(44)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_TITLE
                paragraph.alignment = PP_ALIGN.LEFT

        # Agregar contenido
        if content:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # Procesar contenido línea por línea
            lines = content.split('\n')
            for line_idx, line in enumerate(lines):
                if line_idx > 0:
                    content_frame.add_paragraph()

                if line.startswith('###'):
                    # Subtítulo (###)
                    p = content_frame.paragraphs[-1]
                    p.text = line.lstrip('#').strip()
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(18)
                        run.font.bold = True
                        run.font.color.rgb = COLOR_SUBTITLE
                elif line.startswith('##'):
                    # Encabezado (##)
                    p = content_frame.paragraphs[-1]
                    p.text = line.lstrip('#').strip()
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(24)
                        run.font.bold = True
                        run.font.color.rgb = COLOR_TITLE
                elif line.startswith('- '):
                    # Bullet point
                    p = content_frame.paragraphs[-1]
                    p.text = line[2:].strip()
                    p.level = 1
                    for run in p.runs:
                        run.font.size = Pt(14)
                        run.font.color.rgb = COLOR_TEXT
                elif line.startswith('  - '):
                    # Sub-bullet point
                    p = content_frame.paragraphs[-1]
                    p.text = line[4:].strip()
                    p.level = 2
                    for run in p.runs:
                        run.font.size = Pt(12)
                        run.font.color.rgb = COLOR_TEXT
                elif line.startswith('```'):
                    # Código (ignorar delimitadores)
                    continue
                elif line.strip():
                    # Texto normal
                    p = content_frame.paragraphs[-1]
                    p.text = line.strip()
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(13)
                        run.font.color.rgb = COLOR_TEXT

        # Agregar número de slide en pie de página
        footer_box = slide.shapes.add_textbox(Inches(9), Inches(7.2), Inches(0.8), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"Slide {i}"
        for paragraph in footer_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(180, 180, 180)
            paragraph.alignment = PP_ALIGN.RIGHT

    return prs


def main():
    """Función principal"""
    slides_file = Path("ESTACION-5-SLIDES.md")
    output_file = Path("ESTACION-5-SLIDES.pptx")

    if not PPTX_AVAILABLE:
        print("[ERROR] python-pptx no esta instalado.")
        print("\n[INSTALL] Instalalo con:")
        print("   pip install python-pptx")
        print("\nAlternativamente, usa pandoc:")
        print("   pandoc ESTACION-5-SLIDES.md -o ESTACION-5-SLIDES.pptx")
        return False

    if not slides_file.exists():
        print("[ERROR] No se encontro " + str(slides_file))
        return False

    print("[READ] Leyendo " + str(slides_file) + "...")
    markdown_content = slides_file.read_text(encoding='utf-8')

    print("[PARSE] Parseando slides...")
    slides_data = parse_markdown_slides(markdown_content)
    print("[OK] " + str(len(slides_data)) + " slides encontrados")

    print("[CREATE] Creando presentacion PowerPoint...")
    prs = create_pptx_presentation(slides_data)

    print("[SAVE] Guardando " + str(output_file) + "...")
    prs.save(str(output_file))

    print("\n[SUCCESS] Presentacion creada exitosamente!")
    print("[FILE] " + str(output_file.absolute()))
    print("[SLIDES] " + str(len(slides_data)))
    print("[SIZE] " + str(output_file.stat().st_size / 1024) + " KB")
    print("\n[NEXT]")
    print("[1] Abre " + str(output_file) + " en PowerPoint")
    print("[2] Agrega tu logo/colores corporativos")
    print("[3] Personaliza fuentes y layout")
    print("[4] Guarda como PDF si necesitas compartir")

    return True


if __name__ == "__main__":
    success = main()
    exit(0 if success else 1)
